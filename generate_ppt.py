#!/usr/bin/env python3
"""
MCPFlow 本科毕业答辩 PPT — 15页版本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

# ========== 颜色定义 ==========
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)
DEEPER    = RGBColor(0x06, 0x0F, 0x1E)
CARD_BG   = RGBColor(0x1B, 0x2E, 0x4A)
ACCENT    = RGBColor(0x00, 0x96, 0xD6)
LIGHT_B   = RGBColor(0x48, 0xCA, 0xE4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xC8, 0xD6, 0xE5)
ORANGE    = RGBColor(0xFF, 0xA0, 0x50)
GREEN     = RGBColor(0x4E, 0xC9, 0x8F)
BORDER    = RGBColor(0x2A, 0x4A, 0x6B)
PURPLE    = RGBColor(0xA0, 0x80, 0xFF)
SALMON    = RGBColor(0xF0, 0x80, 0x60)
LIGHT_FILL= RGBColor(0x50, 0x80, 0xF0)
DARK_CARD = RGBColor(0x20, 0x35, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 15

# ========== 工具函数 ==========
def bg(slide, c1=DARK_BG, c2=None):
    if c2:
        fill = slide.background.fill; fill.gradient(); fill.gradient_angle = 135.0
        fill.gradient_stops[0].color.rgb = c1; fill.gradient_stops[1].color.rgb = c2
    else:
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c1

def rect(s, l, t, w, h, fill=CARD_BG, bc=None, r=0.05):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0] = r; return sh

def tb(s, l, t, w, h, text, fs=12, c=WHITE, bold=False, al=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    tx = s.shapes.add_textbox(l, t, w, h); tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = fn; p.alignment = al; return tx

def mtb(s, l, t, w, h, lines, ds=12, dc=WHITE, al=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(l, t, w, h); tf = tx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln.get('text',''); p.font.size = Pt(ln.get('size',ds))
        p.font.color.rgb = ln.get('color',dc); p.font.bold = ln.get('bold',False)
        p.font.name = ln.get('font','Microsoft YaHei'); p.alignment = al
        if 'sp' in ln: p.space_after = Pt(ln['sp'])
    return tx

def dbox(s, l, t, w, h, text, fill=ACCENT, tc=WHITE, fs=10, bold=True):
    sh = rect(s, l, t, w, h, fill=fill)
    tf = sh.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(fs); p.font.color.rgb = tc
    p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    return sh

def line(s, x1, y1, x2, y2, c=BORDER, w=2):
    cn = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = c; cn.line.width = Pt(w); return cn

def aline(s, l, t, w):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

def circ(s, l, t, sz, fill=ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()

def pn(s, n):
    tb(s, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35), f"{n}/{TOTAL}", fs=9, c=GRAY, al=PP_ALIGN.RIGHT)

def header(s, title, num):
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), fill=ACCENT)
    tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5), title, fs=28, c=WHITE, bold=True)
    aline(s, Inches(0.8), Inches(0.85), Inches(1.5))
    pn(s, num)

def arrow_right(s, x, y):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.18), Inches(0.25))
    a.fill.solid(); a.fill.fore_color.rgb = BORDER; a.line.fill.background()

# ====================================================================
# S1: 封面
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DEEPER, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill=ACCENT)
rect(s, Inches(1.5), Inches(1.5), Inches(0.04), Inches(2.5), fill=ACCENT)
dbox(s, Inches(1.8), Inches(1.5), Inches(1.8), Inches(0.38), "MCPFlow", fill=ACCENT, fs=14)
tb(s, Inches(1.8), Inches(2.15), Inches(10.8), Inches(0.8),
   "基于 MCP 协议的多智能体工作流编排平台的设计与实现", fs=30, c=WHITE, bold=True)
rect(s, Inches(1.8), Inches(3.2), Inches(3.0), Inches(0.03), fill=ACCENT)
mtb(s, Inches(1.8), Inches(3.5), Inches(5.0), Inches(3.0), [
    {'text':'答辩人：郑银杰','size':16,'c':GRAY,'sp':10},
    {'text':'专　业：计算机科学与技术','size':16,'c':GRAY,'sp':10},
    {'text':'学　院：计算机学院（人工智能学院）','size':16,'c':GRAY,'sp':10},
    {'text':'指导教师：唐菀 教授','size':16,'c':GRAY,'sp':10},
    {'text':'答辩日期：2026年5月','size':16,'c':GRAY},
])
for i in range(3): rect(s, Inches(10.0), Inches(1.6)+Inches(1.2*i), Inches(2.5), Inches(0.005), fill=RGBColor(0x1B,0x3A,0x5C))
circ(s, Inches(11.5), Inches(6.0), Inches(1.2), fill=RGBColor(0x0A,0x20,0x3E))
circ(s, Inches(12.2), Inches(6.5), Inches(0.6), fill=RGBColor(0x14,0x30,0x50))
rect(s, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), fill=ACCENT)

# ====================================================================
# S2: 目录
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DEEPER, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), fill=ACCENT)
tb(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5), "目　录", fs=32, c=WHITE, bold=True)
aline(s, Inches(0.8), Inches(0.9), Inches(1.5))

toc = [
    ("01", "研究背景与项目目标", ACCENT),
    ("02", "技术选型与系统架构", GREEN),
    ("03", "DAG 工作流引擎设计", ORANGE),
    ("04", "Agent 与 MCP 工具调用", PURPLE),
    ("05", "系统功能模块展示", LIGHT_B),
    ("06", "项目运行效果展示", SALMON),
    ("07", "测试结果与项目总结", GREEN),
]
for i, (num, title, color) in enumerate(toc):
    y = Inches(1.3) + Inches(0.75 * i)
    circ(s, Inches(1.5), y + Inches(0.05), Inches(0.4), fill=color)
    tb(s, Inches(1.6), y + Inches(0.08), Inches(0.2), Inches(0.35), num, fs=12, c=WHITE, bold=True, al=PP_ALIGN.CENTER)
    tb(s, Inches(2.2), y + Inches(0.08), Inches(6.0), Inches(0.4), title, fs=20, c=WHITE, bold=True)
    if i < len(toc)-1:
        rect(s, Inches(2.2), y + Inches(0.55), Inches(8.5), Inches(0.01), fill=BORDER)

pn(s, 2)

# ====================================================================
# S3: 研究背景
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "研究背景", 3)

# 三大背景卡片
cards_data = [
    ("LLM 与 Agent 爆发", "大语言模型能力快速提升，AI Agent 从单一问答走向多步骤自主任务执行，对工具调用和流程编排提出更高要求", ACCENT),
    ("MCP 协议标准化", "Model Context Protocol 提供统一的工具接入规范，使 Agent 能够跨平台发现和调用外部工具，打破能力边界", GREEN),
    ("编排工具缺位", "现有工作流平台多为 SaaS 服务，部署受限；缺少面向多智能体协作的开源、轻量级编排解决方案", ORANGE),
]
for i, (title, desc, color) in enumerate(cards_data):
    x = Inches(0.8) + Inches(4.0 * i)
    rect(s, x, Inches(1.3), Inches(3.6), Inches(3.0), fill=CARD_BG, bc=BORDER)
    rect(s, x, Inches(1.3), Inches(3.6), Inches(0.05), fill=color)
    tb(s, x+Inches(0.25), Inches(1.55), Inches(3.1), Inches(0.4), title, fs=16, c=WHITE, bold=True)
    tb(s, x+Inches(0.25), Inches(2.1), Inches(3.1), Inches(2.0), desc, fs=13, c=GRAY)

# 核心痛点总结
rect(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.6), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.4),
   "核心痛点：复杂 AI 任务需多步骤、多工具协作 → 缺少标准化、可视化的多智能体编排平台", fs=14, c=WHITE, bold=True, al=PP_ALIGN.CENTER)

# MCP 协议简介
rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(5.7), Inches(6.0), Inches(0.35), "🔌 MCP (Model Context Protocol)", fs=16, c=ACCENT, bold=True)
mtb(s, Inches(1.1), Inches(6.15), Inches(11.0), Inches(0.8), [
    {'text':'• MCP 是 Anthropic 提出的开放协议，定义了大模型与外部工具/数据源之间的标准化交互方式','size':12,'c':GRAY,'sp':8},
    {'text':'• 类似"AI 应用的 USB-C 接口"：一次接入 MCP Server，所有兼容 Agent 均可使用其工具','size':12,'c':GRAY,'sp':8},
    {'text':'• 本项目将 MCP 协议集成到工作流引擎中，使 Agent 节点能自动发现并调用外部 MCP Server','size':12,'c':GRAY},
])

pn(s, 3)

# ====================================================================
# S4: 项目目标
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "项目目标", 4)

# 三大目标
goals = [
    ("低代码", "通过可视化拖拽编排工作流" + chr(10) + "降低 AI 流程构建门槛", GREEN),
    ("可扩展", "支持 8 种节点类型" + chr(10) + "灵活接入 MCP 外部工具", ACCENT),
    ("易部署", "前后端分离 + 容器化" + chr(10) + "Docker Compose 一键启动", ORANGE),
]
for i, (title, desc, color) in enumerate(goals):
    x = Inches(0.8) + Inches(4.0 * i)
    rect(s, x, Inches(1.3), Inches(3.6), Inches(3.5), fill=CARD_BG, bc=color if i==0 else BORDER)
    dbox(s, x+Inches(0.5), Inches(1.6), Inches(2.6), Inches(0.6), title, fill=color, fs=18)
    tb(s, x+Inches(0.25), Inches(2.5), Inches(3.1), Inches(1.8), desc, fs=13, c=GRAY, al=PP_ALIGN.CENTER)

# 核心功能列表
rect(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.1), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(5.2), Inches(5.0), Inches(0.35), "📋 项目核心功能", fs=16, c=ACCENT, bold=True)
feats = [
    ("工作流可视化编排", "拖拽节点、连线、配置与保存", ACCENT),
    ("DAG 工作流执行引擎", "拓扑排序、条件分支、变量传递", GREEN),
    ("Agent 节点工具调用", "MCP 工具发现 + Function Calling 循环", ORANGE),
    ("执行监控与日志", "SSE 实时推送、步骤详情记录", PURPLE),
    ("容器化部署", "Docker Compose 一键启动全栈", SALMON),
    ("MCP Server 管理", "注册、连接检测、工具发现", LIGHT_B),
]
for i, (title, desc, color) in enumerate(feats):
    row, col = i//3, i%3
    x = Inches(1.2) + Inches(3.8 * col)
    y = Inches(5.65) + Inches(0.7 * row)
    tb(s, x, y, Inches(3.5), Inches(0.3), f"▸ {title}", fs=12, c=color, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.3), Inches(3.3), Inches(0.3), desc, fs=10, c=GRAY)

pn(s, 4)

# ====================================================================
# S5: 技术选型
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "技术选型", 5)

stacks = [
    ("🖥 前端", [
        ("React", "组件化 UI 框架"),
        ("TypeScript", "类型安全"),
        ("XYFlow", "流程可视化"),
        ("Ant Design", "UI 组件库"),
        ("Vite", "构建工具"),
    ], ACCENT),
    ("⚙️ 后端", [
        ("Go", "高性能编译语言"),
        ("Gin", "HTTP Web 框架"),
        ("mcp-go", "MCP 协议实现"),
        ("Goja", "JS 沙箱执行"),
        ("expr-lang", "表达式求值"),
    ], GREEN),
    ("🗄 数据与部署", [
        ("MongoDB", "文档型数据库"),
        ("Docker", "容器化"),
        ("Nginx", "反向代理"),
        ("SSE", "实时推送"),
        ("OpenAI API", "LLM 接入"),
    ], ORANGE),
]
for i, (title, items, color) in enumerate(stacks):
    x = Inches(0.8) + Inches(4.0 * i)
    rect(s, x, Inches(1.2), Inches(3.6), Inches(4.2), fill=CARD_BG, bc=color)
    tb(s, x+Inches(0.2), Inches(1.35), Inches(3.2), Inches(0.4), title, fs=18, c=color, bold=True)
    for j, (tech, desc) in enumerate(items):
        y = Inches(2.0) + Inches(0.6 * j)
        dbox(s, x+Inches(0.25), y, Inches(1.3), Inches(0.4), tech, fill=color, fs=9)
        tb(s, x+Inches(1.7), y+Inches(0.05), Inches(1.7), Inches(0.35), desc, fs=10, c=GRAY)

# 协议与能力
rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(5.8), Inches(3.0), Inches(0.35), "🔌 协议与能力：", fs=14, c=ACCENT, bold=True)
tb(s, Inches(3.4), Inches(5.8), Inches(8.8), Inches(0.35),
   "MCP  ·  Streamable HTTP  ·  SSE  ·  OpenAI Chat Completions API  ·  Function Calling", fs=13, c=GRAY)

pn(s, 5)

# ====================================================================
# S6: 系统总体架构
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "系统总体架构设计", 6)

# 前端
rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.3), fill=RGBColor(0x15,0x30,0x4A), bc=ACCENT)
tb(s, Inches(1.0), Inches(1.25), Inches(3.0), Inches(0.3), "🖥 前端展示层", fs=14, c=ACCENT, bold=True)
for i, (t, d) in enumerate([("React+TS", "可视化编辑器"),("XYFlow","拖拽连线"),("Ant Design","管理后台"),("Axios+SSE","流式监控")]):
    x = Inches(1.2)+Inches(2.8*i); rect(s, x, Inches(1.6), Inches(2.5), Inches(0.7), fill=CARD_BG, bc=BORDER)
    tb(s, x+Inches(0.15), Inches(1.63), Inches(2.2), Inches(0.28), t, fs=11, c=WHITE, bold=True)
    tb(s, x+Inches(0.15), Inches(1.9), Inches(2.2), Inches(0.3), d, fs=9, c=GRAY)

line(s, Inches(6.5), Inches(2.5), Inches(6.5), Inches(2.9), c=ACCENT, w=2)

# 后端
rect(s, Inches(0.8), Inches(2.9), Inches(7.5), Inches(2.2), fill=RGBColor(0x15,0x30,0x4A), bc=GREEN)
tb(s, Inches(1.0), Inches(2.95), Inches(3.0), Inches(0.3), "⚙️ 后端核心层 (Go + Gin)", fs=14, c=GREEN, bold=True)
for i, (t, d) in enumerate([("REST API","工作流 CRUD"+chr(10)+"MCP管理"),("DAG引擎","拓扑排序"+chr(10)+"条件分支"),("MCP集成","工具发现"+chr(10)+"FuncCall"),("SSE推送","流式状态"+chr(10)+"实时日志")]):
    x = Inches(1.2)+Inches(1.85*i); rect(s, x, Inches(3.35), Inches(1.65), Inches(1.5), fill=CARD_BG, bc=BORDER)
    tb(s, x+Inches(0.1), Inches(3.4), Inches(1.45), Inches(0.3), t, fs=10, c=GREEN, bold=True)
    tb(s, x+Inches(0.1), Inches(3.7), Inches(1.45), Inches(1.0), d, fs=9, c=GRAY)

# 数据库
rect(s, Inches(0.8), Inches(5.3), Inches(7.5), Inches(0.9), fill=RGBColor(0x15,0x30,0x4A), bc=ORANGE)
tb(s, Inches(1.0), Inches(5.35), Inches(3.0), Inches(0.3), "🗄 数据持久层", fs=14, c=ORANGE, bold=True)
dbox(s, Inches(1.2), Inches(5.7), Inches(6.8), Inches(0.4), "MongoDB — 工作流 / 执行记录 / MCP Server / LLM Provider 配置", fill=CARD_BG, tc=GRAY, fs=10)

# 外部
rect(s, Inches(8.8), Inches(2.9), Inches(3.7), Inches(3.3), fill=RGBColor(0x18,0x2A,0x3F), bc=RGBColor(0x50,0x40,0x80))
tb(s, Inches(9.0), Inches(2.95), Inches(3.5), Inches(0.3), "🔌 外部服务", fs=14, c=PURPLE, bold=True)
for i, (t, d) in enumerate([("LLM Provider","OpenAI 兼容"+chr(10)+"Chat Completions"),("MCP Server","工具注册"+chr(10)+"HTTP / SSE")]):
    y = Inches(3.5)+Inches(1.25*i); rect(s, Inches(9.1), y, Inches(3.2), Inches(1.0), fill=CARD_BG, bc=BORDER)
    tb(s, Inches(9.3), y+Inches(0.05), Inches(2.8), Inches(0.28), t, fs=11, c=PURPLE, bold=True)
    tb(s, Inches(9.3), y+Inches(0.32), Inches(2.8), Inches(0.55), d, fs=9, c=GRAY)
line(s, Inches(8.8), Inches(4.1), Inches(8.3), Inches(4.1), w=1)
line(s, Inches(8.8), Inches(5.3), Inches(8.3), Inches(5.3), w=1)

pn(s, 6)

# ====================================================================
# S7: DAG 工作流引擎 — 模型
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "核心设计一：DAG 工作流引擎 —— 工作流模型", 7)

tb(s, Inches(0.8), Inches(1.2), Inches(6.0), Inches(0.4), "📐 工作流 = DAG (有向无环图)", fs=20, c=ACCENT, bold=True)

# 模型组件
for i, (title, items, color) in enumerate([
    ("节点 (Node)", "Start / End / LLM / Agent" + chr(10) + "Condition / Code / HTTP / Email" + chr(10) + "每种节点有独立的配置 Schema", ACCENT),
    ("边 (Edge)", "定义节点间的执行顺序" + chr(10) + "支持条件分支：true / false 两条路径" + chr(10) + "基于表达式结果动态路由", GREEN),
    ("执行上下文", "Key-Value 存储结构" + chr(10) + "保存每个节点的输出结果" + chr(10) + "后续节点通过 {{node_id.output}} 引用", ORANGE),
]):
    x = Inches(0.8) + Inches(4.0 * i)
    rect(s, x, Inches(1.85), Inches(3.6), Inches(3.0), fill=CARD_BG, bc=BORDER)
    rect(s, x, Inches(1.85), Inches(3.6), Inches(0.05), fill=color)
    tb(s, x+Inches(0.2), Inches(2.05), Inches(3.2), Inches(0.35), title, fs=16, c=color, bold=True)
    tb(s, x+Inches(0.2), Inches(2.55), Inches(3.2), Inches(2.0), items, fs=12, c=GRAY)

# 下方：模板变量示例
rect(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(2.0), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(5.3), Inches(8.0), Inches(0.35), "🔗 模板变量机制", fs=16, c=ACCENT, bold=True)
tb(s, Inches(1.1), Inches(5.75), Inches(11.0), Inches(0.3),
   "示例：LLM 节点的 System Prompt 可写为：", fs=12, c=GRAY)
tb(s, Inches(1.1), Inches(6.05), Inches(11.0), Inches(0.3),
   '"根据用户输入 {{start_node.user_input}}，参考搜索引擎结果 {{search_agent.result}}，生成最终回答"', fs=12, c=LIGHT_B, bold=True)
tb(s, Inches(1.1), Inches(6.45), Inches(11.0), Inches(0.5),
   "→ 执行引擎会在运行时将 {{node_id.output}} 替换为对应节点的实际输出值，实现节点间数据流动", fs=12, c=GRAY)

pn(s, 7)

# ====================================================================
# S8: DAG 执行流程
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "核心设计一：DAG 工作流引擎 —— 执行流程", 8)

# 流程步骤
steps = [
    ("1. 拓扑排序", "解析 DAG 图" + chr(10) + "确定执行顺序", GREEN),
    ("2. 初始化上下文", "创建 KV 存储" + chr(10) + "注入用户输入", ACCENT),
    ("3. 按序执行节点", "根据节点类型" + chr(10) + "调用对应 Handler", ORANGE),
    ("4. 条件分支判断", "评估表达式结果" + chr(10) + "选择下游路径", PURPLE),
    ("5. 保存节点输出", "存入执行上下文" + chr(10) + "供后续节点引用", LIGHT_B),
    ("6. 记录日志", "状态、耗时、输出" + chr(10) + "SSE 实时推送", SALMON),
]

for i, (title, desc, color) in enumerate(steps):
    row = i // 3; col = i % 3
    x = Inches(0.8) + Inches(4.0 * col)
    y = Inches(1.3) + Inches(2.5 * row)
    rect(s, x, y, Inches(3.6), Inches(2.1), fill=CARD_BG, bc=color)
    dbox(s, x+Inches(0.4), y+Inches(0.15), Inches(2.8), Inches(0.5), title, fill=color, fs=13)
    tb(s, x+Inches(0.3), y+Inches(0.85), Inches(3.0), Inches(1.1), desc, fs=12, c=GRAY, al=PP_ALIGN.CENTER)

pn(s, 8)

# ====================================================================
# S9: Agent 与 MCP 工具调用
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "核心设计二：Agent 与 MCP 工具调用", 9)

tb(s, Inches(0.8), Inches(1.15), Inches(6.0), Inches(0.35), "🤖 Agent 节点执行流程", fs=18, c=ACCENT, bold=True)

pw = Inches(1.5); ph = Inches(1.1); ag = Inches(0.22)
steps9 = [
    ("连接 MCP"+chr(10)+"Server", ACCENT),
    ("发现工具"+chr(10)+"ListTools", GREEN),
    ("Function"+chr(10)+"Calling 转换", LIGHT_FILL),
    ("LLM 决策"+chr(10)+"工具调用", ORANGE),
    ("执行 MCP"+chr(10)+"callTool", PURPLE),
    ("返回结果"+chr(10)+"给 LLM", SALMON),
    ("生成最终"+chr(10)+"回答", ACCENT),
]
tw = len(steps9)*(pw+ag)-ag; sx = (Inches(13.333)-tw)/2
for i, (text, color) in enumerate(steps9):
    x = sx + i*(pw+ag); y = Inches(1.7)
    circ(s, x+Inches(0.5), y-Inches(0.28), Inches(0.28), fill=color)
    tb(s, x+Inches(0.53), y-Inches(0.26), Inches(0.22), Inches(0.22), str(i+1), fs=9, c=WHITE, bold=True, al=PP_ALIGN.CENTER)
    rect(s, x, y+Inches(0.1), pw, ph, fill=CARD_BG, bc=color)
    tb(s, x+Inches(0.06), y+Inches(0.25), pw-Inches(0.12), ph-Inches(0.3), text, fs=10, c=WHITE, bold=True, al=PP_ALIGN.CENTER)
    if i < len(steps9)-1: arrow_right(s, x+pw+Inches(0.02), y+Inches(0.5))

# 工具调用循环
tb(s, Inches(0.8), Inches(3.3), Inches(6.0), Inches(0.35), "⚡ 工具调用循环 (Tool-Use Loop)", fs=18, c=ACCENT, bold=True)
rect(s, Inches(0.8), Inches(3.75), Inches(5.8), Inches(1.5), fill=CARD_BG, bc=BORDER)
loop_items = [
    "1. LLM 分析任务 → 决定调用工具及参数",
    "2. 系统执行 callTool → 获得工具返回结果",
    "3. 结果返回 LLM → 判断继续/结束",
    "4. 循环直到 LLM 生成最终回答",
]
for i, item in enumerate(loop_items):
    tb(s, Inches(1.1), Inches(3.85)+Inches(0.3*i), Inches(5.3), Inches(0.28), item, fs=11, c=GRAY)

# MCP 兼容性
tb(s, Inches(7.2), Inches(3.3), Inches(6.0), Inches(0.35), "🔌 MCP 协议兼容性", fs=18, c=ACCENT, bold=True)
rect(s, Inches(7.2), Inches(3.75), Inches(5.3), Inches(1.5), fill=CARD_BG, bc=BORDER)
for i, item in enumerate([
    "✅ Streamable HTTP（MCP 推荐）",
    "✅ SSE 自动降级（兼容旧版）",
    "✅ 启动时连接检测 + 工具发现",
    "✅ OpenAI Chat Completions API",
]):
    tb(s, Inches(7.4), Inches(3.85)+Inches(0.3*i), Inches(4.8), Inches(0.28), item, fs=11, c=GREEN)

# Agent 监控
tb(s, Inches(0.8), Inches(5.55), Inches(6.0), Inches(0.35), "📊 Agent 执行监控", fs=18, c=ACCENT, bold=True)
rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), fill=CARD_BG, bc=BORDER)
for i, item in enumerate(["记录每步工具调用详情","展示工具名称与参数","记录工具返回结果","统计执行耗时","SSE 实时推送"]):
    tb(s, Inches(1.2)+Inches(2.3*i), Inches(6.1), Inches(2.0), Inches(0.5), "▸ "+item, fs=10, c=GRAY)

pn(s, 9)

# ====================================================================
# S10: MCP 协议集成细节
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "MCP 协议集成细节", 10)

# 左侧：MCP Server 管理
rect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.0), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(1.3), Inches(5.0), Inches(0.35), "🔧 MCP Server 管理", fs=16, c=ACCENT, bold=True)
for i, item in enumerate([
    "注册：填写 Server 名称、URL、传输方式",
    "连接检测：启动时自动 ping 测试连通性",
    "工具发现：调用 tools/list 获取工具列表",
    "工具调用：通过 tools/call 执行具体工具",
    "支持 Streamable HTTP 与 SSE 双传输模式",
]):
    tb(s, Inches(1.1), Inches(1.8)+Inches(0.4*i), Inches(5.3), Inches(0.35), f"▸ {item}", fs=12, c=GRAY)

# 右侧：Agent 调用流程
rect(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(3.0), fill=CARD_BG, bc=BORDER)
tb(s, Inches(7.5), Inches(1.3), Inches(5.0), Inches(0.35), "🔄 Agent 调用 MCP 流程", fs=16, c=ACCENT, bold=True)
for i, item in enumerate([
    "1. Agent 节点启动，获取关联的 MCP Server 列表",
    "2. 并行调用各 Server 的 tools/list",
    "3. 合并工具列表，转换为 OpenAI Tool 格式",
    "4. 发送给 LLM，LLM 返回 tool_calls",
    "5. 执行 tool_calls，结果回传 LLM",
    "6. 循环直到 LLM 返回最终文本回答",
]):
    tb(s, Inches(7.5), Inches(1.8)+Inches(0.4*i), Inches(4.8), Inches(0.35), item, fs=11, c=GRAY)

# 底部：传输模式对比
rect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(4.6), Inches(6.0), Inches(0.35), "📡 传输模式对比", fs=16, c=ACCENT, bold=True)

modes = [
    ("Streamable HTTP", "MCP 协议推荐方式", "✅ 双向流式通信" + chr(10) + "✅ 连接复用" + chr(10) + "✅ 性能更优", ACCENT),
    ("SSE (Server-Sent Events)", "兼容旧版 Server", "✅ 单向流式推送" + chr(10) + "⚠️ 需额外 HTTP 请求" + chr(10) + "✅ 兼容性好", GREEN),
]
for i, (name, desc, features, color) in enumerate(modes):
    x = Inches(1.2) + Inches(5.6 * i)
    rect(s, x, Inches(5.1), Inches(5.0), Inches(1.6), fill=DARK_CARD, bc=color)
    tb(s, x+Inches(0.2), Inches(5.2), Inches(4.5), Inches(0.3), name, fs=14, c=color, bold=True)
    tb(s, x+Inches(0.2), Inches(5.5), Inches(4.5), Inches(0.3), desc, fs=10, c=GRAY)
    tb(s, x+Inches(0.2), Inches(5.85), Inches(4.5), Inches(0.75), features, fs=11, c=GRAY)

pn(s, 10)

# ====================================================================
# S11: 系统功能模块展示
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "系统功能模块展示", 11)

modules = [
    ("📋 工作流管理","创建、编辑、删除"+chr(10)+"列表展示与搜索",ACCENT),
    ("🎨 可视化编辑器","节点拖拽 + 连线"+chr(10)+"所见即所得编排",GREEN),
    ("⚙️ 节点配置","8 种节点类型"+chr(10)+"独立配置面板",ORANGE),
    ("🔌 MCP Server","注册 + 连接检测"+chr(10)+"工具发现与管理",LIGHT_B),
    ("🤖 LLM Provider","模型提供商"+chr(10)+"API Key 管理",PURPLE),
    ("📊 执行记录","历史执行"+chr(10)+"详情与日志",SALMON),
    ("🧪 Agent Playground","Agent 调试"+chr(10)+"实时对话测试",ACCENT),
    ("📈 统计看板","Dashboard"+chr(10)+"概览与统计",GREEN),
]
cw=Inches(2.75); ch=Inches(2.45); mx=Inches(0.8); my=Inches(1.2); gx=Inches(0.25); gy=Inches(0.25)
for i, (title, desc, color) in enumerate(modules):
    r, c = i//4, i%4; x = mx+c*(cw+gx); y = my+r*(ch+gy)
    rect(s, x, y, cw, ch, fill=CARD_BG, bc=BORDER)
    rect(s, x, y, cw, Inches(0.04), fill=color)
    tb(s, x+Inches(0.2), y+Inches(0.2), cw-Inches(0.4), Inches(0.42), title, fs=14, c=WHITE, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.7), cw-Inches(0.4), Inches(1.5), desc, fs=11, c=GRAY)

pn(s, 11)

# ====================================================================
# S12: 项目展示（上）
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "项目运行效果展示（上）", 12)

screens = [
    ("工作流管理列表","可视化编辑器","节点配置面板"),
    ("执行结果展示","SSE 实时执行状态","工作流编辑画布"),
]
tw_s=Inches(3.7); th=Inches(2.5); gx_s=Inches(0.35); gy_s=Inches(0.3)
for row in range(2):
    for col in range(3):
        x = Inches(0.8) + col*(tw_s+gx_s); y = Inches(1.2) + row*(th+gy_s)
        rect(s, x, y, tw_s, th, fill=RGBColor(0x12,0x22,0x38), bc=BORDER)
        tb(s, x+Inches(0.15), y+Inches(0.12), tw_s-Inches(0.3), Inches(0.28),
           "📷 "+screens[row][col], fs=11, c=ACCENT, bold=True)
        tb(s, x+Inches(0.3), y+Inches(0.7), tw_s-Inches(0.6), Inches(1.2),
           "[ 系统截图："+screens[row][col]+" ]", fs=12, c=RGBColor(0x4A,0x6A,0x8A), al=PP_ALIGN.CENTER)

pn(s, 12)

# ====================================================================
# S13: 项目展示（下）
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "项目运行效果展示（下）", 13)

screens2 = [
    ("Agent Playground","执行日志详情","Docker 部署"),
    ("MCP Server 管理","LLM Provider 管理","数据库管理"),
]
for row in range(2):
    for col in range(3):
        x = Inches(0.8) + col*(tw_s+gx_s); y = Inches(1.2) + row*(th+gy_s)
        rect(s, x, y, tw_s, th, fill=RGBColor(0x12,0x22,0x38), bc=BORDER)
        tb(s, x+Inches(0.15), y+Inches(0.12), tw_s-Inches(0.3), Inches(0.28),
           "📷 "+screens2[row][col], fs=11, c=ACCENT, bold=True)
        tb(s, x+Inches(0.3), y+Inches(0.7), tw_s-Inches(0.6), Inches(1.2),
           "[ 系统截图："+screens2[row][col]+" ]", fs=12, c=RGBColor(0x4A,0x6A,0x8A), al=PP_ALIGN.CENTER)

tb(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.35),
   "✅ 系统已实现从编排 → 配置 → 执行 → 监控 → 日志查看的完整闭环", fs=14, c=GREEN, bold=True, al=PP_ALIGN.CENTER)

pn(s, 13)

# ====================================================================
# S14: 测试结果与项目总结
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, "测试结果与项目总结", 14)

# 测试结果
rect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(3.8), fill=CARD_BG, bc=BORDER)
tb(s, Inches(1.1), Inches(1.3), Inches(5.0), Inches(0.35), "🧪 功能测试结果", fs=16, c=ACCENT, bold=True)
for i, (item, detail) in enumerate([
    ("工作流管理","创建、编辑、删除功能正常，DAG 校验通过"),
    ("多类型节点","8 种节点类型均可正常添加、配置与执行"),
    ("Agent 工具调用","通过 MCP 协议发现工具并完成调用"),
    ("SSE 实时推送","实时展示节点执行状态与输出内容"),
    ("Docker 部署","docker-compose up 一键启动全栈服务"),
]):
    y = Inches(1.8)+Inches(0.6*i)
    tb(s, Inches(1.1), y, Inches(5.3), Inches(0.25), "✅ "+item, fs=12, c=GREEN, bold=True)
    tb(s, Inches(1.1), y+Inches(0.25), Inches(5.3), Inches(0.28), detail, fs=10, c=GRAY)

# 总结
rect(s, Inches(7.2), Inches(1.2), Inches(5.3), Inches(3.8), fill=CARD_BG, bc=BORDER)
tb(s, Inches(7.5), Inches(1.3), Inches(5.0), Inches(0.35), "📝 项目总结", fs=16, c=ACCENT, bold=True)
mtb(s, Inches(7.5), Inches(1.8), Inches(4.8), Inches(2.5), [
    {'text':'本项目设计并实现了 MCPFlow', 'size':13,'c':WHITE,'bold':True,'sp':10},
    {'text':'一个面向多智能体协作的低代码、可视化工作流编排平台。', 'size':12,'c':GRAY,'sp':12},
    {'text':'通过集成 MCP 协议，实现了 Agent 对外部工具的标准化接入与调用。', 'size':12,'c':GRAY,'sp':12},
    {'text':'采用 DAG 引擎 + 模板变量机制，支持复杂任务的有序编排与数据传递。', 'size':12,'c':GRAY,'sp':12},
    {'text':'前后端分离 + 容器化部署，开箱即用。', 'size':12,'c':GRAY,'sp':12},
    {'text':'有效降低了复杂 AI 自动化流程的构建门槛。', 'size':12,'c':GRAY},
])

# 展望
rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.9), fill=CARD_BG, bc=ORANGE)
tb(s, Inches(1.1), Inches(5.4), Inches(6.0), Inches(0.35), "🔮 未来展望", fs=16, c=ORANGE, bold=True)
outlook = [
    ("并行执行","多个独立节点并行调度，提升效率"),
    ("权限管理","多用户角色与权限控制"),
    ("版本管理","工作流版本控制与回滚"),
    ("节点扩展","更多节点类型与插件机制"),
    ("定时调度","Cron 定时触发工作流"),
    ("分布式","支持多实例部署与负载均衡"),
]
for i, (title, detail) in enumerate(outlook):
    row, col = i//3, i%3
    x = Inches(1.2) + Inches(3.8 * col)
    y = Inches(5.85) + Inches(0.6 * row)
    tb(s, x, y, Inches(3.5), Inches(0.25), "▸ "+title, fs=11, c=ORANGE, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.25), Inches(3.3), Inches(0.25), detail, fs=10, c=GRAY)

pn(s, 14)

# ====================================================================
# S15: 致谢
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DEEPER, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), fill=ACCENT)
rect(s, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), fill=ACCENT)
tb(s, Inches(0), Inches(1.5), Inches(13.333), Inches(1.2), "致  谢", fs=48, c=WHITE, bold=True, al=PP_ALIGN.CENTER)
aline(s, Inches(5.5), Inches(2.8), Inches(2.333))
mtb(s, Inches(2.0), Inches(3.2), Inches(9.333), Inches(3.0), [
    {'text':'感谢唐菀教授在选题、系统设计和论文撰写过程中的悉心指导','size':16,'c':GRAY,'sp':20},
    {'text':'感谢计算机学院（人工智能学院）各位老师四年来的培养与帮助','size':16,'c':GRAY,'sp':20},
    {'text':'感谢同学们在学习和生活中的支持与陪伴','size':16,'c':GRAY,'sp':20},
    {'text':'感谢家人一直以来的理解与鼓励','size':16,'c':GRAY,'sp':20},
], al=PP_ALIGN.CENTER)
rect(s, Inches(5.0), Inches(5.8), Inches(3.333), Inches(0.02), fill=ACCENT)
tb(s, Inches(0), Inches(6.0), Inches(13.333), Inches(0.6), "恳请各位老师批评指正", fs=18, c=WHITE, bold=True, al=PP_ALIGN.CENTER)
pn(s, 15)

# ========== 保存 ==========
out = '/Users/zhengyinjie/project/mcpflow/MCPFlow_毕业答辩.pptx'
prs.save(out)
print(f"PPT generated: {out}")
print(f"Total slides: {len(prs.slides)}")
